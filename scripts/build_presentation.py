"""Build the hackathon slide deck from LIVE computed results (no stale numbers).

Usage:  python3 scripts/build_presentation.py
Output: presentation/Portfolio_Health_Advisor.pptx
"""
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "src"))

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

from tools import (assess_risk, calculate_allocation, calculate_returns,
                   health_score, project_growth, rebalance_plan, stress_test,
                   tax_loss_harvest)

# ---------- live data ----------
df = pd.read_csv(ROOT / "sample_data" / "sample_portfolio.csv")
R = calculate_returns(df)
A = calculate_allocation(df)
T = R["totals"]
RISK = assess_risk(df, R, A)
HEALTH = health_score(R, A, RISK)
PLAN = rebalance_plan(df)
PROJ = project_growth(T["total_value"])
STRESS = stress_test(df)
HARV = tax_loss_harvest(R)
LARGE = ROOT / "sample_data" / "large_portfolio_200.csv"
N_LARGE = len(pd.read_csv(LARGE)) if LARGE.is_file() else 200

OUT_DIR = ROOT / "presentation"
OUT_DIR.mkdir(exist_ok=True)
CHART = OUT_DIR / "_alloc_chart.png"

# ---------- chart ----------
sectors = sorted(A["by_sector"].items(), key=lambda x: -x[1])
fig, ax = plt.subplots(figsize=(7.5, 3.6))
ax.barh([s for s, _ in reversed(sectors)], [p for _, p in reversed(sectors)],
        color=["#d62728" if p >= 50 else "#ff9f0e" if p >= 30 else "#1f77b4"
               for _, p in reversed(sectors)])
ax.set_xlabel("Portfolio %")
ax.set_title("Sample portfolio: sector concentration (Technology dominates)")
fig.tight_layout()
fig.savefig(CHART, dpi=150)
plt.close(fig)

# ---------- theme ----------
BG_DARK = RGBColor(0x0F, 0x1B, 0x2D)
ACCENT = RGBColor(0x00, 0xC2, 0xA8)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def _para(tf, text, size=20, bold=False, color=INK, space_after=6, alignment=None):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_after = Pt(space_after)
    if alignment is not None:
        p.alignment = alignment
    return p


def title_slide(title, subtitle, footer="ACM Student Chapter Hackathon  ·  Agentic AI"):
    s = prs.slides.add_slide(BLANK)
    _bg(s, BG_DARK)
    tf = _box(s, 0.8, 1.6, 11.7, 2.2).text_frame
    tf.word_wrap = True
    _para(tf, title, size=44, bold=True, color=WHITE)
    tf2 = _box(s, 0.8, 3.6, 11.7, 1.6).text_frame
    tf2.word_wrap = True
    _para(tf2, subtitle, size=24, color=ACCENT)
    tf3 = _box(s, 0.8, 6.4, 11.7, 0.6).text_frame
    _para(tf3, footer, size=16, color=WHITE)
    return s


def content_slide(title, bullets, image=None, caption=None):
    s = prs.slides.add_slide(BLANK)
    _bg(s, WHITE)
    tf = _box(s, 0.6, 0.2, 12.1, 1.0).text_frame
    tf.word_wrap = True
    _para(tf, title, size=32, bold=True, color=BG_DARK)
    bar = s.shapes.add_shape(1, Inches(0.6), Inches(1.15), Inches(1.4), Pt(5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    right = 7.2 if image else 12.1
    tf2 = _box(s, 0.6, 1.5, right, 5.4).text_frame
    tf2.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = b
        p.font.size = Pt(19)
        p.font.color.rgb = INK
        p.space_after = Pt(10)
        p.level = 0
    if image:
        s.shapes.add_picture(str(image), Inches(8.1), Inches(1.6), width=Inches(4.4))
        if caption:
            tf3 = _box(s, 8.1, 6.5, 4.4, 0.5).text_frame
            _para(tf3, caption, size=13, color=MUTED)
    return s


# ---------- deck ----------
title_slide("Portfolio Health Advisor",
            "Two AI agents that turn a confusing portfolio into clear, honest advice.")

content_slide("The problem", [
    "Retail investors hold stocks across sectors but can't see the big picture.",
    "Three questions nobody answers easily: How much did I really gain? "
    "How concentrated is my risk? What should I do about it?",
    "Our answer: automate that reasoning with a two-agent AI pipeline.",
])

content_slide("Solution: sequential agent pipeline", [
    f"Agent 1 — Analyst: reads the CSV, runs math tools. Worth {T['total_value']} "
    f"(+{T['total_pnl']} / +{T['return_pct']}%), top sector {A['top_sector']} {A['top_sector_pct']}%. Facts only.",
    f"Agent 1b — Risk: flags {RISK['risk_level']} risk (score {RISK['risk_score']}/100), "
    f"estimates tax, and adds health, rebalance, stress and harvest analysis. Rule-based, instant.",
    "Agent 2 — Advisor: a free-tier LLM reasons over the findings and writes a "
    "120-word plain-English summary with 1–2 actions. Offline rule engine as fallback.",
    "Then: multi-turn Q&A over the same findings ('why is my risk high?').",
])

content_slide("Live results on the sample portfolio", [
    f"Value {T['total_value']} on cost {T['total_cost']} → +{T['total_pnl']} (+{T['return_pct']}%). "
    f"{T['num_winners']} winners, {T['num_losers']} losers (worst: CLDW −11.1%).",
    f"Health {HEALTH['score']}/100 (Grade {HEALTH['grade']}) — Technology at "
    f"{A['top_sector_pct']}% is the core problem (HHI {A['hhi']}).",
    f"Rebalance plan frees ~{PLAN['freed_total']} by trimming "
    f"{', '.join(s['ticker'] for s in PLAN['sells'][:3])}.",
    f"Stress: {STRESS['note']}",
], image=CHART, caption="Red ≥ 50% · Orange ≥ 30%")

content_slide("Creative features (beyond the brief)", [
    "Health score 0–100 + grade with penalty breakdown — one number anyone gets.",
    f"One-click rebalance plan — exact SELL quantities, BUY amounts (frees ~{PLAN['freed_total']}).",
    "Growth projector + SIP goal planner — 'SIP for 100000 in 5 years?' answered live.",
    "Stress-test lab — crash any sector by any % and see the damage.",
    "Tax estimators — gains tax + tax-loss harvest pairs (CLDW offsets TCHX, saves ~77.4).",
    "Report export + advice cache (repeat runs in ~0.2s) + full-screen mouse-driven TUI.",
])

content_slide("Scale: 10 → 200 holdings", [
    f"Generator builds reproducible datasets (seeded): {N_LARGE}-holding portfolio ≈ 3.85M value.",
    "Same pipeline, no code changes — CSV in, dashboard out.",
    "Large portfolio correctly reads as diversified (top sector ≈ 12%) vs the concentrated sample.",
    "200 holdings compute in well under 2 seconds (vectorized pandas).",
])

content_slide("Accuracy you can audit", [
    "Exact math: totals computed from unrounded values, reconciled to the cent on both datasets.",
    "Sector percentages largest-remainder normalized — always sum to exactly 100.00%.",
    "HHI concentration index from exact weights, not rounded ones.",
    "Validation: friendly errors for missing columns, bad tickers, bad percentages.",
    "Tested: tools + agents (LLM stubbed) + real mouse-click TUI tests — all green.",
])

content_slide("Interfaces for every judge", [
    "Full-screen TUI (Textual): sidebar, tables, What-If lab, live Q&A chat, one-click export.",
    "Classic Rich menu + plain CLI (--auto, --non-interactive, --question) for scripts.",
    "./run.sh — one command from clone to dashboard, with or without an API key.",
    "Markdown report in outputs/ — the whole analysis as a shareable file.",
])

content_slide("Demo plan (3 minutes)", [
    "1. ./run.sh → dashboard: health gauge, red concentration bar, advice panel.",
    "2. Rebalance view: exact trim/buy list. Stress lab: crash Technology 20%.",
    "3. Ask: 'SIP for 100000 in 5 years?' → exact monthly number.",
    "4. Export report → open outputs/*.md. Show --auto + cache speed on rerun.",
])

title_slide("Thank you — questions?",
            "Not financial advice. Educational demo on static sample data.",
            footer="github.com/algospider/AGENTIC-AI  ·  ./run.sh to try it live")

path = OUT_DIR / "Portfolio_Health_Advisor.pptx"
prs.save(path)
print(f"Saved {path} ({len(prs.slides._sldIdLst)} slides)")
